#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Mon Apr 24 09:08:02 2023

@author: matias.otthgmail.com
"""


import streamlit as st
import pandas as pd
import plotly.express as px
import numpy as np
from datetime import time
from pptx import Presentation
from io import BytesIO
from datetime import date
from pptx.dml.color import RGBColor
from pptx.util import Inches
import os


def add_image(slide, image, left, top, width):
    slide.shapes.add_picture(image, left=left, top=top, width=width)


st.set_page_config(layout="wide")
st.title('MONITOR ECONÓMICO CPP USS')
today = date.today()

st.header('Generar presentación con información económica')


#path="/Users/matias.otthgmail.com/Desktop/Monitor_Economico/"
data=pd.read_parquet("datafull.parquet")
data1=data[data["CATEGORIA"]=="ACTIVIDAD ECONOMICA"]
data2=data[data["CATEGORIA"]=="INFLACION"]
data3=data[data["CATEGORIA"]=="MERCADO LABORAL"]
data4=data[data["CATEGORIA"]=="CUENTAS CORRIENTES"]


extremos_1=[data1["PERIODO"].iloc[0].to_pydatetime(),data1["PERIODO"].iloc[-1].to_pydatetime()]
extremos_2=[data2["PERIODO"].iloc[0].to_pydatetime(),data2["PERIODO"].iloc[-1].to_pydatetime()]
extremos_3=[data3["PERIODO"].iloc[0].to_pydatetime(),data3["PERIODO"].iloc[-1].to_pydatetime()]
extremos_4=[data4["PERIODO"].iloc[0].to_pydatetime(),data4["PERIODO"].iloc[-1].to_pydatetime()]


options = [ "ACTIVIDAD ECONÓMICA","INFLACIÓN","MERCADO LABORAL","CUENTAS CORRIENTES"]
user_input = st.multiselect(label='Selecciones la serie a utilizar', options=options)

dic_options={"ACTIVIDAD ECONÓMICA":["IMACEC","CRECIMIENTO ECONÓMICO"],
             "INFLACIÓN":["YoY","MENSUAL"],
             "MERCADO LABORAL":["DESOCUPACIÓN","OCUPACIÓN Y PARTICIPACIÓN"],
             "CUENTAS CORRIENTES":["TOTAL","DESAGREGADAS"]
             }




if options[0] in user_input:
    serie=options[0]
    st.subheader(serie)
    user_input_1 = st.multiselect(label='Selecciones la serie a utilizar', options=dic_options[serie])
    appointment_1 = st.slider(
        "Seleccione el rango de fechas para la serie " + serie,
        value=(extremos_1[0],extremos_1[1]),
        format="YYYY/MM/DD")

    
if options[1] in user_input:
    serie=options[1]
    st.subheader(serie)
    user_input_2 = st.multiselect(label='Selecciones la serie a utilizar', options=dic_options[serie])
    appointment_2 = st.slider(
        "Seleccione el rango de fechas para la serie " + serie,
        value=(extremos_2[0],extremos_2[1]),
        format="YYYY/MM/DD")

      
    
if options[2]in user_input:
    serie=options[2]
    st.subheader(serie)
    user_input_3 = st.multiselect(label='Selecciones la serie a utilizar', options=dic_options[serie])
    appointment_3 = st.slider(
        "Seleccione el rango de fechas para la serie " + serie,
        value=(extremos_3[0],extremos_3[1]),
        format="YYYY/MM/DD")


if options[3] in user_input:
    serie=options[3]
    st.subheader(serie)
    user_input_4 = st.multiselect(label='Selecciones la serie a utilizar', options=dic_options[serie])
    appointment_4 = st.slider(
        "Seleccione el rango de fechas para la serie " + serie,
        value=(extremos_4[0],extremos_4[1]),
        format="YYYY/MM/DD")
    




sub1 = st.checkbox(label='Exportar como presentación')

if sub1:
    genre = st.radio(
        "Seleccionar un formato",
        ('Formato claro', 'Formato oscuro'))

    
    col1, col2 = st.columns(2)
    with col1:
       st.subheader("Formato claro")
       st.image("OPCION1.png")
       
    with col2:
       st.subheader("Formato oscuro")
       st.image("OPCION2.png")
    
   
    
    try:
        data1=data1[(data1["PERIODO"]> appointment_1[0])&(data1["PERIODO"]< appointment_1[1])]
    except:
        pass
    try:
        data2=data2[(data2["PERIODO"]> appointment_2[0])&(data2["PERIODO"]< appointment_2[1])]
    except:
        pass
    try:
        data3=data3[(data3["PERIODO"]> appointment_3[0])&(data3["PERIODO"]< appointment_3[1])]
    except:
        pass
    try:
        data4=data4[(data4["PERIODO"]> appointment_4[0])&(data4["PERIODO"]< appointment_4[1])]
    except:
        pass
    
    #ACTIVIDAD ECONÓMICA
    imacec = px.line(data1[data1["SERIE"]=="1.Imacec"], x="PERIODO", y="VALOR", color="SERIE", template='simple_white')
    # imacec.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')

    componentes_imacec=px.line(data1[data1["CATEGORIA2"]=="IMACEC"], x="PERIODO", y="VALOR", color="SERIE", template='simple_white')
    # componentes_imacec.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')
   
    pib_anual= px.line(data1[data1["SERIE"]=="PIB ANUAL"], x="PERIODO", y="VALOR", color="SERIE", template='simple_white')
    # pib_anual.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')
 
    pib_TRIMESTRAL= px.line(data1[data1["SERIE"].isin(["YoY","Desestacionalizado (Variación Trimestral)"])], x="PERIODO", y="VALOR", color="SERIE", template='simple_white')
    # pib_TRIMESTRAL.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')

    
    #INFLACIÓN
    
    ipcs=["Variación Mensual", "YoY"]
    ipc_a= px.line(data2[data2["SERIE"]==ipcs[1]], x="PERIODO", y="VALOR", color="SERIE", template='simple_white')
    # ipc_a.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')

    ipc_m= px.line(data2[data2["SERIE"]==ipcs[0]], x="PERIODO", y="VALOR", color="SERIE", template='simple_white')
    # ipc_m.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')

    #MERCADO LABORAL
    desocupados=px.line(data3[data3["SERIE"]=="Tasa de desocupación"], x="PERIODO", y="VALOR", color="SERIE", template='simple_white')
    # desocupados.user_input_update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')

    oc_part=px.line(data3[data3["SERIE"].isin(["Tasa de ocupación","Tasa de participación"])], x="PERIODO", y="VALOR", color="SERIE", template='simple_white')
    # oc_part.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')

      
    #CUENTAS CORRIENTES
    cuentas=px.line(data4[(data4["SERIE"]=="TOTAL")&~(data4["CATEGORIA2"]=="Jurídica")], x="PERIODO", y="VALOR", color="SERIE", template='simple_white')
    # cuentas.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')

    cuentas_2=px.line(data4[(data4["SERIE"]=="TOTAL")&(data4["CATEGORIA2"]=="Jurídica")], x="PERIODO", y="VALOR", color="SERIE", template='simple_white')
    # cuentas_2.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')
            
    desagregadas=px.line(data4[~(data4["SERIE"]=="TOTAL")&~(data4["CATEGORIA2"]=="Jurídica")], x="PERIODO", y="VALOR", color="SERIE", template='simple_white')
    # desagregadas.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')

    desagregadas_2=px.line(data4[~(data4["SERIE"]=="TOTAL")&(data4["CATEGORIA2"]=="Jurídica")], x="PERIODO", y="VALOR", color="SERIE", template='simple_white')
    # desagregadas_2.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')

        
    title_1 = st.text_input('Título de la presentación', 'Informe de Actividad Económica')
    submit = st.button(label='GENERAR PRESENTACIÓN')
    
    
    
    if submit and user_input == "":
        st.warning("Selecionar una portada")
    
    elif submit and user_input != "":
        with st.spinner('Preparando... :man_cook_tone1:'):
            prs=Presentation("BASE PRESENTACION USS_STGO-BES.pptx")
            xml_slides = prs.slides._sldIdLst  
            slides = list(xml_slides)
            
            if genre=="Formato claro" :           
                xml_slides.remove(slides[1]) 
            else: 
                xml_slides.remove(slides[0]) 
            
            #GENERAR PRESENTACIÓN
            slide = prs.slides[0]
            title =  slide.shapes.title.text_frame.paragraphs[0]
            title.text = title_1
            title.font.color.rgb = RGBColor(255, 255, 255)  # Color blanco
            title.font.bold = True  # Negrita


            #AGREGAR GRAFICOS
            width = Inches(8)
            left = Inches(5)
            top = Inches(1)
            
        
            imacec.write_image("imacec.png")
            im="imacec.png"
            add_image(prs.slides[2], image=im, left=left, width=width, top=top)
            os.remove("imacec.png")


            componentes_imacec.write_image("com_imacec.png")
            com_im="com_imacec.png"
            add_image(prs.slides[3], image=com_im, left=left, width=width, top=top)
            os.remove("com_imacec.png")
    
    
            pib_anual.write_image("pib_a.png")
            pib_aa="pib_a.png"
            add_image(prs.slides[4], image=pib_aa, left=left, width=width, top=top)
            os.remove("pib_a.png")
    
    
            pib_TRIMESTRAL.write_image("pib_t.png")
            pib_t="pib_t.png"
            add_image(prs.slides[5], image=pib_t, left=left, width=width, top=top)
            os.remove("pib_t.png")
    
            ipc_a.write_image("ipc_a.png")
            ipc_aa="ipc_a.png"
            add_image(prs.slides[7], image=ipc_aa, left=left, width=width, top=top)
            os.remove("ipc_a.png")
        
            ipc_m.write_image("ipc_m.png")
            ipc_mm="ipc_m.png"
            add_image(prs.slides[8], image=ipc_mm, left=left, width=width, top=top)
            os.remove("ipc_m.png")
    
            desocupados.write_image("desocupados.png")
            desocupadoss="desocupados.png"
            add_image(prs.slides[10], image=desocupadoss, left=left, width=width, top=top)
            os.remove("desocupados.png")
 
            oc_part.write_image("oc_part.png")
            oc_partt="oc_part.png"
            add_image(prs.slides[11], image=oc_partt, left=left, width=width, top=top)
            os.remove("oc_part.png")
    
            cuentas.write_image("cuentas.png")
            cuentass="cuentas.png"
            add_image(prs.slides[13], image=cuentass, left=left, width=width, top=top)
            os.remove("cuentas.png")
        
            desagregadas.write_image("desagregadas.png")
            desagregadass="desagregadas.png"
            add_image(prs.slides[14], image=desagregadass, left=left, width=width, top=top)
            os.remove("desagregadas.png")
    
            try:       
                user_input_1=pd.DataFrame(user_input_1)
            except:
                pass
            
            try:       
                user_input_2=pd.DataFrame(user_input_2)
            except:
                pass
            try:       
                user_input_3=pd.DataFrame(user_input_3)
            except:
                pass
            try:       
                user_input_4=pd.DataFrame(user_input_4)
            except:
                pass
            
            
            sacar=[]
            try:
                if "IMACEC" not in user_input_1.values and "CRECIMIENTO ECONÓMICO" not in user_input_1.values:
                    sacar.append(2)

                if "IMACEC" not in user_input_1.values:
                    sacar.append(3)
                    sacar.append(4)
                if "CRECIMIENTO ECONÓMICO" not in user_input_1.values:
                    sacar.append(5)
                    sacar.append(6)
            except:
                pass
            try:
                if "YoY" not in user_input_2.values and "MENSUAL" not in user_input_2.values:
                    sacar.append(7)

                if "YoY" not in user_input_2.values:
                    sacar.append(8)
                if "MENSUAL" not in user_input_2.values:
                    sacar.append(9)           
            except:
                pass  
            
            try:
                if "DESOCUPACIÓN" not in user_input_3.values and "OCUPACIÓN Y PARTICIPACIÓN" not in user_input_3.values:    
                    sacar.append(10)
                if "DESOCUPACIÓN" not in user_input_3.values:
                    sacar.append(11)
                if "OCUPACIÓN Y PARTICIPACIÓN" not in user_input_3.values:
                    sacar.append(12)
            except:
                pass
            
            try:
                if "TOTAL" not in user_input_4.values and "DESAGREGADAS" not in user_input_4.values:
                    sacar.append(13)
                if "TOTAL" not in user_input_4.values:
                    sacar.append(14)
                if "DESAGREGADAS" not in user_input_4.values:
                    sacar.append(15)
            except:
                pass   
            
              
            x=1
            for i in sacar:
                xml_slides = prs.slides._sldIdLst  
                slides = list(xml_slides)
                xml_slides.remove(slides[i-x]) 
                x=x+1
            
            

            #GENERAR ARCHIVO
            filename = 'presentacion_{}_{}.pptx'.format("economia", today)
            binary_output = BytesIO()
            prs.save(binary_output)
            st.download_button(label='Presiona para descargar',
                               data=binary_output.getvalue(),
                               file_name=filename)
else:
    pass





