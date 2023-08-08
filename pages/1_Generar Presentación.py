#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
@author: matias.otthgmail.com
"""

import streamlit as st
import pandas as pd
import plotly.express as px
import numpy as np
from datetime import time
from pptx import Presentation
from io import BytesIO
from datetime import date
from pptx.dml.color import RGBColor
from pptx.util import Inches
import os
from datetime import datetime
from pptx.util import Pt
import calendar

st.set_page_config(layout="wide")
st.image("logo.png")

#st.sidebar.image("ESCUDOUSS_vertical_color.png", use_column_width=True)

def add_image(slide, image, left, top, width):
    slide.shapes.add_picture(image, left=left, top=top, width=width)

nombre_mes_espanol = [
    "enero", "febrero", "marzo", "abril", "mayo", "junio",
    "julio", "agosto", "septiembre", "octubre", "noviembre", "diciembre"
]



def porcentaje(dato, decimas):
    porcentaje_str = str(np.round(dato * 100, decimas)).replace(".", ",") + "%"
    return porcentaje_str

def gen(imacec_des,rango,titulo):
    imacec_des=imacec_des[(imacec_des["PERIODO"]> rango[0])&(imacec_des["PERIODO"]< rango[1])]
    imacec_des = px.line(imacec_des, x="PERIODO", y="VALOR", color="SERIE" ,title='Mi gráfico de línea', 
              labels={'x': 'Eje X', 'y': 'Eje Y'}, 
              template='plotly_white', 
              width=700, height=600)
    imacec_des.update_layout(title={
        'text': titulo,
        'x':0.5,
         'xanchor': 'center',
         'yanchor': 'top' 
          },legend=dict(
            orientation="h",
            yanchor="top",
            y=-0.2,
            xanchor="left",
            x=0.01
        ))
    return imacec_des



st.markdown("<h1 style='text-align: center; color: black;'>MONITOR ECONÓMICO CPP USS</h1>", unsafe_allow_html=True)
st.markdown("<h2 style='text-align: center; color: grey;'>Generar presentación con información económica</h2>", unsafe_allow_html=True)


st.write(' ')

today = date.today()



data=pd.read_parquet("datos_monitor_1.parquet")


data1=data[data["CATEGORIA"]=="ACTIVIDAD ECONOMICA"]
data2=data[data["CATEGORIA"]=="INFLACION"]
data3=data[data["CATEGORIA"]=="MERCADO LABORAL"]
data4=data[data["CATEGORIA"]=="CUENTAS CORRIENTES"]



extremos_1=[data1["PERIODO"].iloc[0].to_pydatetime(),data2["PERIODO"].iloc[-1].to_pydatetime()]
extremos_2=[data2["PERIODO"].iloc[0].to_pydatetime(),data2["PERIODO"].iloc[-1].to_pydatetime()]
extremos_3=[data3["PERIODO"].iloc[0].to_pydatetime(),data3["PERIODO"].iloc[-1].to_pydatetime()]



options = [ "ACTIVIDAD ECONÓMICA","INFLACIÓN","MERCADO LABORAL"]
user_input = st.multiselect(label='Selecciones la serie a utilizar', options=options)

dic_options={"ACTIVIDAD ECONÓMICA":["ACTIVIDAD","COMPONENTES"],
             "INFLACIÓN":["ANUAL"],
             "MERCADO LABORAL":["DESOCUPACIÓN","INFORMALIDAD","REMUNERACIONES"]
             }


submit=st.checkbox(label='Seleccionar todas las categorías')
if submit:
    user_input=["ACTIVIDAD ECONÓMICA","INFLACIÓN","MERCADO LABORAL"]
    


if options[0] in user_input:
    serie=options[0]
    st.subheader(serie)
    user_input_1 = st.multiselect(label='Selecciones la serie a utilizar', options=dic_options[serie])
    appointment_1 = st.slider("Seleccione el rango de fechas para la series. " + serie,
                   value=(extremos_1[0],extremos_1[1]),
                   format="YYYY/MM")

     

if options[1] in user_input:
    serie=options[1]
    st.subheader(serie)
    user_input_2 = st.multiselect(label='Selecciones la serie a utilizar', options=dic_options[serie])
   
    appointment_2 = st.slider(
                            "Seleccione el rango de fechas para la serie.  " + serie,
                            value=(extremos_2[0],extremos_2[1]),
                            format="YYYY/MM")

      
  
if options[2] in user_input:
    serie=options[2]
    st.subheader(serie)
    user_input_3 = st.multiselect(label='Selecciones la serie a utilizar', options=dic_options[serie])
    
    appointment_3 = st.slider(
                        "Seleccione el rango de fechas para la serie.   " + serie,
                        value=(extremos_3[0],extremos_3[1]),
                        format="YYYY/MM")

sub1 = st.checkbox(label='Exportar como presentación')

if sub1:
    genre = st.radio(
        "Seleccionar un formato",
        ('Formato CPP', 'Formato Facultad'))

    col1, col2 = st.columns(2)
    with col1:
       st.subheader("Plantilla Centro Políticas Públicas")
       st.image("cpp.png")

    with col2:
       st.subheader("Plantilla Facultad Economía y Gobierno")
       st.image("facultad.png")

    title_1 = st.text_input('Título de la presentación', 'Informe de Actividad Económica')
    submit = st.button(label='GENERAR PRESENTACIÓN')



    if submit and user_input == "":
        st.warning("Selecionar una portada")

    elif submit and user_input != "":
        with st.spinner('Generando presentación...'):
            


            try:
                data1=data1[(data1["PERIODO"]> appointment_1[0])&(data1["PERIODO"]< appointment_1[1])]
            except:
                pass
            try:
                data2=data2[(data2["PERIODO"]> appointment_2[0])&(data2["PERIODO"]< appointment_2[1])]
            except:
                pass
            try:
                data3=data3[(data3["PERIODO"]> appointment_3[0])&(data3["PERIODO"]< appointment_3[1])]
            except:
                pass
            try:
                data4=data4[(data4["PERIODO"]> appointment_4[0])&(data4["PERIODO"]< appointment_4[1])]
            except:
                pass

            #ACTIVIDAD ECONÓMICA
            def gen(imacec_des,rango,titulo):
                imacec_des=imacec_des[(imacec_des["PERIODO"]>= rango[0])&(imacec_des["PERIODO"]<= rango[1])]
                imacec_des = px.line(imacec_des, x="PERIODO", y="VALOR", color="SERIE" ,title='Mi gráfico de línea', 
                          labels={'x': 'Eje X', 'y': 'Eje Y'}, 
                          template='plotly_white', 
                          width=700, height=600)
                imacec_des.update_yaxes(rangemode="tozero")
                imacec_des.update_layout(title={
                    'text': titulo,
                    'x':0.5,
                     'xanchor': 'center',
                     'yanchor': 'top' 
                      },legend=dict(
                        orientation="h",
                        yanchor="top",
                        y=-0.2,
                        xanchor="left",
                        x=0.01
                    ))
                return imacec_des
            def fechas_2(grafico):
                grafico.update_xaxes(
                    rangeselector=dict(
                        buttons=list([
                            dict(count=3, label="3A", step="year", stepmode="backward"),
                            dict(count=5, label="5A", step="year", stepmode="backward"),
                            dict(count=10, label="10A", step="year", stepmode="backward"),
                            dict(step="all")
                        ])
                    )
                )
                return grafico

            def eje_porcentaje(grafico):
                grafico.layout.yaxis.tickformat = ',.0%'
                return grafico    

            def gen_bar(imacec_des,rango,titulo):
                imacec_des=imacec_des[(imacec_des["PERIODO"] >=rango[0])&(imacec_des["PERIODO"] <= rango[1])]    
                imacec_des = px.bar(imacec_des, x="PERIODO", y="VALOR", color="SERIE", title='Mi gráfico de línea', 
                          labels={'x': 'Eje X', 'y': 'Eje Y'}, 
                          template='plotly_white', 
                          width=700, height=600)
                          # color_discrete_map=rename_dict)
                imacec_des.update_yaxes(rangemode="tozero")
                imacec_des.update_layout(title={
                    'text': titulo,
                    'x':0.5,
                     'xanchor': 'center',
                     'yanchor': 'top' 
                      },legend=dict(
                        orientation="h",
                        yanchor="top",
                        y=-0.2,
                        xanchor="left",
                        x=0.01
                    ))
                return imacec_des

            #SLIDE 1 ACTIVIAD ECONOMICA
            try:
                data11=data1[data1["CATEGORIA2"]=="IMACEC"]    
                imacec_or="Imacec empalmado, serie original (índice 2018=100)"
                imacec_or=data11[data11["NOMBRE_2"]==imacec_or]
                imacec_or["VALOR"]=imacec_or["VALOR"]/imacec_or["VALOR"].shift(12)-1
                uv_imacec_or=imacec_or["VALOR"].iloc[-1]
                FECHA_IMACEC = nombre_mes_espanol[imacec_or["PERIODO"].iloc[-1].month-1]
    
                imacec_or=imacec_or.dropna()
                imacec_or["SERIE"]="Imacec (variación anual)"
                imacec_or_1=gen(imacec_or,appointment_1,"Variación anual del IMACEC")
                imacec_or_1=eje_porcentaje(imacec_or_1)

                data13=data1[data1["CATEGORIA2"]=="PIB"]
                nom="PIB, volumen a precios del año anterior encadenado, referencia 2018 (miles de millones de pesos encadenados)"
                nom=data13[data13["NOMBRE_2"]==nom]
                nom["VALOR"]=nom["VALOR"]/nom["VALOR"].shift(4)-1
                uv_nom=nom["VALOR"].iloc[-1]
                FECHA_PIB= nombre_mes_espanol[nom["PERIODO"].iloc[-1].month-1]
                
                nom=nom.dropna()
                nom["SERIE"]="PIB Trimestral (variación YoY)"
                nom=gen(nom,appointment_1,"Variación Trimestral PIB")
                nom=eje_porcentaje(nom)

            #SLIDE 2 COMPONENTES

                data12=data1[data1["CATEGORIA2"]=="IMACEC - COMPONENTES"]
                data12["VALOR"]=data12["VALOR"]/100
                est="Indicador mensual de actividad económica, Imacec, contribución porcentual respecto de igual periodo del año anterior, referencia 2018"
                est=data12[data12["NOMBRE_1"]==est]
                est["SERIE"]=est["NOMBRE_2"]
                
                
              
                prod_bienes=est[est["NOMBRE_2"].isin(["Minería","Industria","Resto de bienes"])]
                
                uv_mineria=est[est["NOMBRE_2"]=="Minería"]["VALOR"].iloc[-1]              
                uv_indsutria=est[est["NOMBRE_2"]=="Industria"]["VALOR"].iloc[-1]
                uv_sericios=est[est["NOMBRE_2"]=="Resto de bienes"]["VALOR"].iloc[-1]
                   
                prod_bienes=gen_bar(prod_bienes,appointment_1,"Componentes producción de bienes")
                prod_bienes=eje_porcentaje(prod_bienes)


             
                componentes=est[est["NOMBRE_2"].isin(["Producción de bienes","Comercio","Servicios"])]
                componentes=gen_bar(componentes,appointment_1,"Componentes principales IMACEC")
                componentes=eje_porcentaje(componentes)
                
            except:
                pass
            #SLIDE 3 
            try:
                data2=data[data["CATEGORIA"]=="INFLACION"]
                data2["VALOR"]=data2["VALOR"]/100
                anu="IPC, IPC sin volátiles e IPC volátiles, variación anual, información empalmada"
                inf_anu=data2[(data2["NOMBRE_1"]==anu)&(data2["NOMBRE_2"]=="IPC General")]
                uv_inf=inf_anu["VALOR"].iloc[-1]
                FECHA_IPC = nombre_mes_espanol[inf_anu["PERIODO"].iloc[-1].month-1]          
                inf_anu["SERIE"]=inf_anu["NOMBRE_2"]

                inf_anu=gen(inf_anu,appointment_2,"Variación anual porcentual IPC")
                inf_anu=eje_porcentaje(inf_anu)

                com_anu=data2[(data2["NOMBRE_1"]==anu)&~(data2["NOMBRE_2"]=="IPC General")]
                com_anu["SERIE"]=com_anu["NOMBRE_2"]
                comp_2=com_anu[~com_anu["SERIE"].isin(["IPC sin volátiles","IPC volátil"])]
               
                uv_servicios=com_anu[com_anu['NOMBRE_2'] == "IPC Servicios sin volátiles"]['VALOR'].iloc[-1]
                uv_bienes=com_anu[com_anu['NOMBRE_2'] == "IPC Bienes sin volátiles"]['VALOR'].iloc[-1]
                uv_alimentos=com_anu[com_anu['NOMBRE_2'] == "IPC Alimentos volátiles"]['VALOR'].iloc[-1]
                uv_energia=com_anu[com_anu['NOMBRE_2'] == "IPC Energía volátiles"]['VALOR'].iloc[-1]
                uv_volatiles=com_anu[com_anu['NOMBRE_2'] == "IPC Resto de volátiles"]['VALOR'].iloc[-1]
                  
                comp_2.loc[comp_2['NOMBRE_2'] == "IPC Servicios sin volátiles", 'VALOR'] = comp_2.loc[comp_2['NOMBRE_2'] == "IPC Servicios sin volátiles", 'VALOR']*0.384
                comp_2.loc[comp_2['NOMBRE_2'] == "IPC Bienes sin volátiles", 'VALOR'] = comp_2.loc[comp_2['NOMBRE_2'] == "IPC Bienes sin volátiles", 'VALOR']*0.267
                comp_2.loc[comp_2['NOMBRE_2'] == "IPC Alimentos volátiles", 'VALOR'] = comp_2.loc[comp_2['NOMBRE_2'] =="IPC Alimentos volátiles" , 'VALOR']*0.101
                comp_2.loc[comp_2['NOMBRE_2'] == "IPC Energía volátiles", 'VALOR'] = comp_2.loc[comp_2['NOMBRE_2'] =="IPC Energía volátiles", 'VALOR']*0.075
                comp_2.loc[comp_2['NOMBRE_2'] == "IPC Resto de volátiles", 'VALOR'] = comp_2.loc[comp_2['NOMBRE_2'] == "IPC Resto de volátiles", 'VALOR']*0.172
                
                comp_2=gen_bar(comp_2,appointment_2,"Componentes secundarias IPC")
                comp_2=eje_porcentaje(comp_2)  
                inf_anu_=inf_anu1[(inf_anu1["PERIODO"] >= appointment_2[0])&(inf_anu1["PERIODO"]<=appointment_2[1])]
                comp_2.add_trace(px.line(inf_anu_, x='PERIODO', y='VALOR', color="SERIE").data[0])

                
            except:
                pass

                #SLIDE 4
            try:
                data3=data[data["CATEGORIA"]=="MERCADO LABORAL"]
                data3["VALOR"]=data3["VALOR"]/100  
                emp_tasas_nac=data3[(data3["CATEGORIA2"]=="EMPLEO - TASAS")&(data3["CATEGORIA3"]=="Nacional")]
                oc=emp_tasas_nac[emp_tasas_nac["NOMBRE_1"]=="Tasa de desocupación Nacional"]
                
                ult_oc=oc["VALOR"].iloc[-1]
                FECHA_INE = nombre_mes_espanol[oc["PERIODO"].iloc[-1].month-1]  
       
                oc["SERIE"]=oc["NOMBRE_2"]
        
                oc=gen(oc,appointment_3,"Tasa de desocupación")
                oc=eje_porcentaje(oc)

                emp_tasas_nac2=data3[(data3["CATEGORIA2"]=="EMPLEO - TASAS")&~(data3["CATEGORIA3"]=="Nacional")]
                oc2=emp_tasas_nac2[emp_tasas_nac2["NOMBRE_1"].isin(["Tasa de desocupación H","Tasa de desocupación M"])]
                
                ult_oc_h=emp_tasas_nac2[emp_tasas_nac2["NOMBRE_1"]=="Tasa de desocupación H"]["VALOR"].iloc[-1]
                ult_oc_m=emp_tasas_nac2[emp_tasas_nac2["NOMBRE_1"]=="Tasa de desocupación M"]["VALOR"].iloc[-1]
   
                oc2["SERIE"]=oc2["NOMBRE_2"]
                oc2=gen(oc2,appointment_3,"Tasas de desocupación por género")
                oc2=eje_porcentaje(oc2)

                informalidad=data3[(data3["CATEGORIA2"]=="INFORMALIDAD")&(data3["NOMBRE_1"]=="Tasa de informalidad (AS)")]

                
                ult_informalidad=informalidad["VALOR"].iloc[-1]
                informalidad["SERIE"]=informalidad["NOMBRE_2"]

                informalidad=informalidad.sort_values(by="PERIODO")
                informalidad=gen(informalidad,appointment_3,"Tasa de Informalidad")
                informalidad=eje_porcentaje(informalidad)
                
                informalidad2=data3[(data3["CATEGORIA2"]=="INFORMALIDAD")&~(data3["NOMBRE_1"]=="Tasa de informalidad (AS)")]
                
                ult_informalidad_h=informalidad2[informalidad2["NOMBRE_1"]=="Tasa de informalidad (H)"]["VALOR"].iloc[-1]    
                ult_informalidad_m=informalidad2[informalidad2["NOMBRE_1"]=="Tasa de informalidad (M)"]["VALOR"].iloc[-1]    
                
                informalidad2["SERIE"]=informalidad2["NOMBRE_2"]
                informalidad2=informalidad2.sort_values(by="PERIODO")
                informalidad2=gen(informalidad2,appointment_3,"Tasas de Informalidad por género")

                informalidad2=eje_porcentaje(informalidad2)

            #SLIDE 5

                ind_rem_men_r=data[(data["CATEGORIA2"]=="INDICE DE REMUNERACIONES")&(data["CATEGORIA3"]=="REAL")]
                ind_rem_men_n=data[(data["CATEGORIA2"]=="INDICE DE REMUNERACIONES")&(data["CATEGORIA3"]=="NOMINAL")]
                ind_rem_men_r["SERIE"]="Variación real Y/Y"     
                ind_rem_men_n["SERIE"]="Variación nominal Y/Y"    
                ind_rem_men_r["VALOR"]=ind_rem_men_r["VALOR"]/ind_rem_men_r["VALOR"].shift(12)-1
                ind_rem_men_r=ind_rem_men_r.dropna()
                ind_rem_men_n["VALOR"]=ind_rem_men_n["VALOR"]/ind_rem_men_n["VALOR"].shift(12)-1
                ind_rem_men_n=ind_rem_men_n.dropna()
                ind_rem_men_r=ind_rem_men_r.sort_values(by="PERIODO")
                ind_rem_men_n=ind_rem_men_n.sort_values(by="PERIODO")
                ult_remuneraciones=ind_rem_men_r["VALOR"].iloc[-1]

                ind_rem_men_r=gen(ind_rem_men_r,appointment_3,"Variación anual Índice de remuneraciones [real]")
                ind_rem_men_r=eje_porcentaje(ind_rem_men_r)

                ind_rem_men_n=gen(ind_rem_men_n,appointment_3,"Variación anual Índice de remuneraciones [nom]")
                ind_rem_men_n=eje_porcentaje(ind_rem_men_n)
            except:
                pass
            
            

            if genre=="Formato CPP" :           
                prs=Presentation("plantilla-cpp.pptx")
                xml_slides = prs.slides._sldIdLst  
                slides = list(xml_slides)
            else: 
                prs=Presentation("plantilla-facultad.pptx")
                xml_slides = prs.slides._sldIdLst  
                slides = list(xml_slides)

            #GENERAR PRESENTACIÓN
            slide = prs.slides[0]
            title =  slide.shapes.title.text_frame.paragraphs[0]
            title.text = title_1
            title.font.color.rgb = RGBColor(255, 255, 255)  # Color blanco
            
            title.font.bold = True  # Negrita


            #AGREGAR GRAFICOS
            width = Inches(5.75)
            leftd = Inches(0.5)
            lefti = Inches(7)
            top= Inches(2)
                 
                #SLIDE 1 ACTIVIAD ECONOMICA        
            try:
                #SLIDE 1 ACTIVIAD ECONOMICA
                imacec_or_1.write_image("imacec.png")
                im1="imacec.png"
                add_image(prs.slides[2], image=im1, left=leftd, width=width, top=top)
                os.remove("imacec.png")

                nom.write_image("nom.png")
                im2="nom.png"
                add_image(prs.slides[2], image=im2, left=lefti, width=width, top=top)
                os.remove("nom.png")    
            except:
                pass 
            
            try:
                #SLIDE 2 COMPONENTES
                prod_bienes.write_image("prod.png")
                im3="prod.png"
                add_image(prs.slides[3], image=im3, left=leftd, width=width, top=top)
                os.remove("prod.png")

                componentes.write_image("comp.png")
                im4="comp.png"
                add_image(prs.slides[3], image=im4, left=lefti, width=width, top=top)
                os.remove("comp.png")
                
            except:
                pass
            try:
                slide2 = prs.slides[2]
                texto = "El IMACEC de "+FECHA_IMACEC+" anotó una variación anual del "  +porcentaje(uv_imacec_or,2)+", el PIB trimestral de "+FECHA_PIB+" anotó una variación anual del "+porcentaje(uv_nom,1)+"."
                title_2 = slide2.shapes.title.text_frame.paragraphs[0]
                title_2.text = texto
                title_2.font.color.rgb = RGBColor(0, 0, 0)  # Color blanco
                title_2.font.name = "Calibri" 
                title_2.font.size = Pt(18)
  

                slide2 = prs.slides[3]
                texto = "La componentes de Minería, Industria y servicios en "+FECHA_IMACEC+" obtuvieron una variación anual del " +porcentaje(uv_mineria,1) +", " + porcentaje(uv_indsutria,1)+", " +porcentaje(uv_sericios,1)+ "."
                title_2 = slide2.shapes.title.text_frame.paragraphs[0]
                title_2.text = texto
                title_2.font.color.rgb = RGBColor(0, 0, 0)  # Color blanco
                title_2.font.name = "Calibri" 
                title_2.font.size = Pt(18)
            except:
                pass
            
            try:
                #SLIDE 3 
                inf_anu.write_image("inf_a.png")
                im5="inf_a.png"
                add_image(prs.slides[5], image=im5, left=leftd, width=width, top=top)
                os.remove("inf_a.png")
                comp_2.write_image("comp_2.png")
                im6="comp_2.png"
                add_image(prs.slides[5], image=im6, left=lefti, width=width, top=top)
                os.remove("comp_2.png")
            except:
                pass
            
            try:
                mayor = np.max(np.abs([uv_servicios, uv_bienes, uv_alimentos, uv_energia, uv_volatiles]))
                etiqueta = np.argmax(np.abs([uv_servicios, uv_bienes, uv_alimentos, uv_energia, uv_volatiles]))
                etiquetas={0:"Servicios no volátiles",
                        1:"Bienes no volátiles",
                        2:"Alimentos",
                        3:"Energía",
                        4:"Resto de volátiles"}
                
                slide2 = prs.slides[5]
                texto = "La variación anual del IPC en "+FECHA_IPC+" alcanzó un "+ porcentaje(uv_inf,1)+". La mayor componente fue " +etiquetas[etiqueta] + " con un " + porcentaje(mayor,1)+"."
                title_2 = slide2.shapes.title.text_frame.paragraphs[0]
                title_2.text = texto
                title_2.font.color.rgb = RGBColor(0, 0, 0)  # Color blanco
                title_2.font.name = "Calibri" 
                title_2.font.size = Pt(18)

            except:
                pass        
            
            try:
                #SLIDE 4
                oc.write_image("desocup.png")
                im7="desocup.png"
                add_image(prs.slides[7], image=im7, left=leftd, width=width, top=top)
                os.remove("desocup.png")

                oc2.write_image("desoc2.png")
                im8="desoc2.png"
                add_image(prs.slides[7], image=im8, left=lefti, width=width, top=top)
                os.remove("desoc2.png")

                informalidad.write_image("infor.png")
                im9="infor.png"
                add_image(prs.slides[8], image=im9, left=leftd, width=width, top=top)
                os.remove("infor.png")
                informalidad2.write_image("infor2.png")
                im10="infor2.png"
                add_image(prs.slides[8], image=im10, left=lefti, width=width, top=top)
                os.remove("infor2.png")

                #SLIDE 5
                ind_rem_men_r.write_image("real.png")
                im11="real.png"
                add_image(prs.slides[9], image=im11, left=leftd, width=width, top=top)
                os.remove("real.png")

                ind_rem_men_n.write_image("nominal.png")
                im12="nominal.png"
                add_image(prs.slides[9], image=im12, left=lefti, width=width, top=top)
                os.remove("nominal.png")
            except:
                pass

            try:
                slide2 = prs.slides[7]
                texto = "La desocupación laboral en "+FECHA_INE+" alcanzó el "+ porcentaje(ult_oc,1) + ",mujeres se encuentran en el "+ porcentaje(ult_oc_m,1) +" y hombres en"+ porcentaje(ult_oc_h,1)+"."
                title_2 = slide2.shapes.title.text_frame.paragraphs[0]
                title_2.text = texto
                title_2.font.color.rgb = RGBColor(0, 0, 0)  # Color blanco
                title_2.font.name = "Calibri" 
                title_2.font.size = Pt(18)

                slide3 = prs.slides[8]
                texto2 = "La Informalidad laboral en "+FECHA_INE+" alcanzó el "+ porcentaje(ult_informalidad,1) + ",mujeres se encuentran en el "+ porcentaje(ult_informalidad_m,1) +" y hombres en"+ porcentaje(ult_informalidad_h,1)+"."
           
                title_3 = slide3.shapes.title.text_frame.paragraphs[0]
                title_3.text = texto2
                title_3.font.color.rgb = RGBColor(0, 0, 0)  # Color blanco
                title_3.font.name = "Calibri" 
                title_3.font.size = Pt(18)
                
                slide4 = prs.slides[9]
                texto3 = "Los salarios reales en "+FECHA_INE+" anotaron una variación anual del "+ porcentaje(ult_remuneraciones,1) +"."
                title_4 = slide4.shapes.title.text_frame.paragraphs[0]
                title_4.text = texto3
                title_4.font.color.rgb = RGBColor(0, 0, 0)  # Color blanco
                title_4.font.name = "Calibri" 
                title_4.font.size = Pt(18)
            except:
                pass

            
            try:       
                user_input_1=pd.DataFrame(user_input_1)
            except:
                pass

            try:       
                user_input_2=pd.DataFrame(user_input_2)
            except:
                pass
            try:       
                user_input_3=pd.DataFrame(user_input_3)
            except:
                pass
       #     try:       
       #         user_input_4=pd.DataFrame(user_input_4)
        #    except:
        #        pass


            sacar=[]

            try:
                if "ACTIVIDAD" not in user_input_1.values:
                    sacar.append(3)
                if "COMPONENTES" not in user_input_1.values:
                    sacar.append(4)
            except: 
                sacar.append(2)
                sacar.append(3)
                sacar.append(4)
            
            try:
                if "ANUAL" not in user_input_2.values:
                    sacar.append(6)
            except:
                sacar.append(5)
                sacar.append(6)
  


            try:
                if "DESOCUPACIÓN" not in user_input_3.values:
                    sacar.append(8)
                if "INFORMALIDAD" not in user_input_3.values:
                    sacar.append(9)
                if "REMUNERACIONES" not in user_input_3.values:
                    sacar.append(10)

            except:
                sacar.append(7)
                sacar.append(8)
                sacar.append(9) 
                sacar.append(10) 

            x=1
            for i in sacar:
                xml_slides = prs.slides._sldIdLst  
                slides = list(xml_slides)
                xml_slides.remove(slides[i-x]) 
                x=x+1




            #GENERAR ARCHIVO
            filename = 'presentación_{}_{}.pptx'.format("Economía", today)
            binary_output = BytesIO()
            prs.save(binary_output)
            st.download_button(label='Descargar',
                               data=binary_output.getvalue(),
                               file_name=filename)
else:
    pass


